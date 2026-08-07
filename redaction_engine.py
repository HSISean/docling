"""
redaction_engine.py

Architecture
------------
Docling (docling.document_converter.DocumentConverter) is used as the
universal parsing/detection layer: it can open PDFs, Word docs, PowerPoint
decks, HTML, images and more, and hand back clean plain text regardless of
source format. We run the user's detection rules (built-in PII patterns,
custom keywords, custom regex) against that docling-extracted text to decide
*what* needs to be redacted and to show a live preview/count in the UI.

Docling does not write documents back out, so the actual redaction (removing
or blacking-out the matched content in the *original* file) is performed by
format-native libraries that understand each container format:
    - PDF   -> PyMuPDF (fitz): true redaction (annotation + apply_redactions,
              which strips the underlying text/glyphs, not just a visual box)
    - DOCX  -> python-docx: run-level text replacement
    - PPTX  -> python-pptx: shape/text-frame run-level replacement
    - TXT/MD -> plain text substitution
    - anything else docling can read -> exported to a redacted .txt with a
      note, since we cannot safely rewrite an unknown binary container.
"""

from __future__ import annotations

import re
import traceback
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Optional

from patterns import DOCLING_PREVIEW_EXTENSIONS


@dataclass
class Rule:
    label: str
    pattern: re.Pattern


LogFn = Optional[Callable[[str], None]]


class RedactionEngine:
    def __init__(
        self,
        rules: list[Rule],
        style: str = "box",          # "box" (block glyphs) or "label" ([REDACTED])
        log_fn: LogFn = None,
    ):
        self.rules = rules
        self.style = style
        self._log = log_fn or (lambda msg: None)
        self._converter = None

    # ------------------------------------------------------------------ #
    # docling (detection / preview)
    # ------------------------------------------------------------------ #
    @property
    def converter(self):
        if self._converter is None:
            self._log("Initializing docling document converter...")
            from docling.document_converter import DocumentConverter
            self._converter = DocumentConverter()
        return self._converter

    def extract_text_with_docling(self, path: Path) -> str:
        """Universal text extraction used for detection/preview across formats."""
        result = self.converter.convert(str(path))
        return result.document.export_to_markdown()

    def find_matches(self, text: str) -> dict[str, int]:
        """Return {label: match_count} for the current rule set against text."""
        counts: dict[str, int] = {}
        for rule in self.rules:
            found = rule.pattern.findall(text)
            if found:
                counts[rule.label] = counts.get(rule.label, 0) + len(found)
        return counts

    def scan(self, path: Path) -> dict[str, int]:
        """Run docling extraction + detection, used for the pre-redaction preview."""
        ext = path.suffix.lower()
        if ext not in DOCLING_PREVIEW_EXTENSIONS:
            return {}
        try:
            text = self.extract_text_with_docling(path)
        except Exception as exc:  # noqa: BLE001
            self._log(f"  (docling preview failed for {path.name}: {exc})")
            return {}
        return self.find_matches(text)

    # ------------------------------------------------------------------ #
    # redaction dispatch
    # ------------------------------------------------------------------ #
    def redact(self, path: Path, output_dir: Path) -> Path:
        output_dir.mkdir(parents=True, exist_ok=True)
        ext = path.suffix.lower()
        try:
            if ext == ".pdf":
                return self._redact_pdf(path, output_dir)
            if ext == ".docx":
                return self._redact_docx(path, output_dir)
            if ext == ".pptx":
                return self._redact_pptx(path, output_dir)
            if ext in (".txt", ".md"):
                return self._redact_text(path, output_dir)
            return self._redact_fallback(path, output_dir)
        except Exception as exc:  # noqa: BLE001
            self._log(f"  ERROR redacting {path.name}: {exc}")
            self._log(traceback.format_exc(limit=2))
            raise

    # ------------------------------------------------------------------ #
    # string-level substitution shared by text-bearing formats
    # ------------------------------------------------------------------ #
    def _redact_string(self, text: str) -> tuple[str, int]:
        count = 0

        def make_sub():
            def _sub(m: re.Match) -> str:
                nonlocal count
                count += 1
                if self.style == "label":
                    return "[REDACTED]"
                return "\u2588" * max(len(m.group(0)), 1)
            return _sub

        for rule in self.rules:
            text = rule.pattern.sub(make_sub(), text)
        return text, count

    def _collect_tokens(self, text: str) -> set[str]:
        """All distinct raw matches across all rules, longest-first so that
        e.g. a full credit-card match is redacted before shorter overlapping
        substrings."""
        tokens: set[str] = set()
        for rule in self.rules:
            for m in rule.pattern.finditer(text):
                token = m.group(0)
                if token and token.strip():
                    tokens.add(token)
        return tokens

    # ------------------------------------------------------------------ #
    # PDF — true redaction via PyMuPDF
    # ------------------------------------------------------------------ #
    def _redact_pdf(self, path: Path, output_dir: Path) -> Path:
        import fitz  # PyMuPDF

        doc = fitz.open(str(path))
        total = 0
        for page_index, page in enumerate(doc):
            page_text = page.get_text()
            tokens = self._collect_tokens(page_text)
            # Redact longer tokens first to avoid leaving remnants of a
            # shorter overlapping match (e.g. part of a credit card number).
            for token in sorted(tokens, key=len, reverse=True):
                areas = page.search_for(token)
                for rect in areas:
                    page.add_redact_annot(rect, fill=(0, 0, 0))
                    total += 1
            if tokens:
                page.apply_redactions()
        out_path = self._unique_path(output_dir, path.stem, "_redacted.pdf")
        doc.save(str(out_path), garbage=4, deflate=True)
        doc.close()
        self._log(f"  PDF: {total} region(s) permanently redacted across {page_index + 1} page(s).")
        return out_path

    # ------------------------------------------------------------------ #
    # DOCX — python-docx run-level replacement
    # ------------------------------------------------------------------ #
    def _redact_docx(self, path: Path, output_dir: Path) -> Path:
        from docx import Document

        doc = Document(str(path))
        total = 0

        def redact_paragraph(paragraph) -> None:
            nonlocal total
            for run in paragraph.runs:
                if not run.text:
                    continue
                new_text, n = self._redact_string(run.text)
                if n:
                    run.text = new_text
                    total += n

        for p in doc.paragraphs:
            redact_paragraph(p)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for p in cell.paragraphs:
                        redact_paragraph(p)

        for section in doc.sections:
            for hdr_ftr in (section.header, section.footer):
                for p in hdr_ftr.paragraphs:
                    redact_paragraph(p)

        out_path = self._unique_path(output_dir, path.stem, "_redacted.docx")
        doc.save(str(out_path))
        self._log(f"  DOCX: {total} match(es) redacted.")
        return out_path

    # ------------------------------------------------------------------ #
    # PPTX — python-pptx shape/run-level replacement
    # ------------------------------------------------------------------ #
    def _redact_pptx(self, path: Path, output_dir: Path) -> Path:
        from pptx import Presentation

        prs = Presentation(str(path))
        total = 0

        def redact_shape(shape) -> None:
            nonlocal total
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if not run.text:
                            continue
                        new_text, n = self._redact_string(run.text)
                        if n:
                            run.text = new_text
                            total += n
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                new_text, n = self._redact_string(run.text)
                                if n:
                                    run.text = new_text
                                    total += n
            if shape.shape_type == 6:  # GROUP
                for sub_shape in shape.shapes:
                    redact_shape(sub_shape)

        for slide in prs.slides:
            for shape in slide.shapes:
                redact_shape(shape)
            if slide.has_notes_slide:
                for shape in slide.notes_slide.shapes:
                    redact_shape(shape)

        out_path = self._unique_path(output_dir, path.stem, "_redacted.pptx")
        prs.save(str(out_path))
        self._log(f"  PPTX: {total} match(es) redacted.")
        return out_path

    # ------------------------------------------------------------------ #
    # Plain text / markdown
    # ------------------------------------------------------------------ #
    def _redact_text(self, path: Path, output_dir: Path) -> Path:
        content = path.read_text(encoding="utf-8", errors="replace")
        new_content, n = self._redact_string(content)
        out_path = self._unique_path(output_dir, path.stem, f"_redacted{path.suffix}")
        out_path.write_text(new_content, encoding="utf-8")
        self._log(f"  TEXT: {n} match(es) redacted.")
        return out_path

    # ------------------------------------------------------------------ #
    # Fallback for other docling-readable formats (xlsx, html, images, etc.)
    # ------------------------------------------------------------------ #
    def _redact_fallback(self, path: Path, output_dir: Path) -> Path:
        self._log(f"  No native writer for '{path.suffix}'; exporting redacted text via docling.")
        text = self.extract_text_with_docling(path)
        new_text, n = self._redact_string(text)
        out_path = self._unique_path(output_dir, path.stem, "_redacted.txt")
        header = (
            f"# Redacted text export of {path.name}\n"
            f"# Original format '{path.suffix}' has no in-place redaction writer;\n"
            f"# this is a docling-extracted, redacted plain-text version.\n\n"
        )
        out_path.write_text(header + new_text, encoding="utf-8")
        self._log(f"  FALLBACK: {n} match(es) redacted in extracted text.")
        return out_path

    # ------------------------------------------------------------------ #
    @staticmethod
    def _unique_path(directory: Path, stem: str, suffix: str) -> Path:
        candidate = directory / f"{stem}{suffix}"
        i = 2
        while candidate.exists():
            candidate = directory / f"{stem}{suffix.rsplit('.', 1)[0]}_{i}.{suffix.rsplit('.', 1)[1]}"
            i += 1
        return candidate


def build_rules(
    builtin_selected: dict[str, str],
    custom_keywords: list[str],
    custom_regex: list[str],
    case_sensitive: bool = False,
    whole_word: bool = False,
) -> list[Rule]:
    """Compile the user's UI selections into a list of Rule objects."""
    flags = 0 if case_sensitive else re.IGNORECASE
    rules: list[Rule] = []

    for label, pattern in builtin_selected.items():
        rules.append(Rule(label, re.compile(pattern, flags)))

    for kw in custom_keywords:
        kw = kw.strip()
        if not kw:
            continue
        escaped = re.escape(kw)
        if whole_word:
            escaped = rf"\b{escaped}\b"
        rules.append(Rule(f'Keyword: "{kw}"', re.compile(escaped, flags)))

    for rx in custom_regex:
        rx = rx.strip()
        if not rx:
            continue
        try:
            rules.append(Rule(f"Custom regex: {rx}", re.compile(rx, flags)))
        except re.error as exc:
            raise ValueError(f"Invalid regex '{rx}': {exc}") from exc

    return rules
